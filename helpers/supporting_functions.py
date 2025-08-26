import requests
import json
import base64
import os
from typing import Dict, List
import pandas as pd
import time
import fitz  # PyMuPDF
from PIL import Image
import io
from pptx import Presentation
from pptx.util import Inches
from googleapiclient.http import MediaIoBaseUpload
from pathlib import Path
from gdrive_module import create_folder, search_files, delete_file, upload_file
import mysql.connector
from collections import defaultdict
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload

import browsercookie

def get_cookies(cookies_path):
    cj = browsercookie.chrome()
    output = []
    for i, cookie in enumerate(cj):
        if ('.google.com' == cookie.domain) or ('lookerstudio.google.com' in cookie.domain):
            cookie_dict = {
                "domain": cookie.domain,
                "expirationDate": int(cookie.expires) if cookie.expires else int(time.time()) + 3600,
                "hostOnly": False,  # browsercookie doesn't provide this directly
                "httpOnly": cookie._rest.get("HttpOnly", False),
                "name": cookie.name,
                "path": cookie.path,
                "sameSite": "unspecified",  # not available in browsercookie
                "secure": cookie.secure,
                "session": False if cookie.expires else True,
                "storeId": "0",  # static unless you retrieve from a specific profile
                "value": cookie.value,
                "id": i
            }
            output.append(cookie_dict)

    with open(cookies_path, "w") as f:
        json.dump(output, f, indent=4)



def pdf_to_pptx_from_bytes(pdf_bytes: bytes, pptx_path: str, scale: float = 1.0) -> str:
    """
    Convert PDF bytes to PPTX format, with optional image scaling.
    Args:
        pdf_bytes: PDF file content in bytes
        pptx_path: Full path to save the PPTX file
        scale: Scale factor for image size (default 1.0 = fit to slide)
    Returns:
        str: Path to the created PPTX file, or empty string if conversion failed
    """
    try:
        # Create output directory if it doesn't exist
        pptx_dir = os.path.dirname(pptx_path)
        os.makedirs(pptx_dir, exist_ok=True)

        # Create a new PowerPoint presentation
        presentation = Presentation()
        # Set slide dimensions to 16:9 (10 inches x 5.625 inches)
        presentation.slide_width = Inches(10)
        presentation.slide_height = Inches(5.625)

        # Open the PDF from bytes using PyMuPDF
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")

        # Determine resampling method for resizing
        if hasattr(Image, 'Resampling'):
            resample_method = Image.Resampling.LANCZOS
        else:
            resample_method = Image.LANCZOS

        for page_num in range(pdf_document.page_count):
            page = pdf_document.load_page(page_num)
            pix = page.get_pixmap()
            img = Image.open(io.BytesIO(pix.tobytes("png")))

            # Resize image if scale != 1.0
            if scale != 1.0:
                new_width = int(img.width * scale)
                new_height = int(img.height * scale)
                img = img.resize((new_width, new_height), resample=resample_method)

            # Save the image to a temporary file
            temp_path = os.path.join(pptx_dir, f"temp_page_{page_num + 1}.png")
            img.save(temp_path, "PNG")

            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide_width = presentation.slide_width
            slide_height = presentation.slide_height

            # Set fixed slide dimensions for 16:9 aspect ratio
            slide_width = Inches(10)  # 10 inches
            slide_height = Inches(5.625)  # 5.625 inches
            
            # Calculate image dimensions to fit the slide
            img_width = slide_width
            img_height = slide_height
            
            # Calculate position to center the image
            left = 0
            top = 0
            
            # Add the image to the slide with exact dimensions
            slide.shapes.add_picture(temp_path, left, top, width=img_width, height=img_height)

            os.remove(temp_path)

        presentation.save(pptx_path)
        print(f"Successfully converted PDF to PPTX: {pptx_path}")
        return pptx_path

    except Exception as e:
        print(f"Error converting PDF to PPTX: {str(e)}")
        return ""


def download_looker_report_bytes(brand_name: str, country_code: str, start_date: str, end_date: str, cookies_file_path: str) -> bytes:
    """
    Download Looker Studio report and return PDF content as bytes.
    Returns PDF bytes if successful, None if failed.
    """
    import json
    import base64
    import requests

    url = "https://lookerstudio.google.com/getPdf?appVersion=20250519_0000"

    # Load the cookies
    with open(cookies_file_path, 'r') as f:
        all_cookies = json.load(f)

    cookie_str = '; '.join([f"{cookie['name']}={cookie['value']}" for cookie in all_cookies])
    xsrf_token = next((cookie['value'] for cookie in all_cookies if cookie['name'] == 'RAP_XSRF_TOKEN'), None)

    if not xsrf_token:
        print("Error: RAP_XSRF_TOKEN not found in cookies file")
        return None
    
    headers = {
        "accept": "application/json, text/plain, */*",
        "content-type": "application/json",
        "encoding": "null",
        "referer": "https://lookerstudio.google.com/reporting/527b536c-e3ab-44ba-918a-9bc7fb639145/page/p_lfee22gzrd/edit",
        "sec-ch-ua": '"Chromium";v="136", "Google Chrome";v="136", "Not.A/Brand";v="99"',
        "sec-ch-ua-mobile": "?1",
        "sec-ch-ua-platform": '"Android"',
        "user-agent": "Mozilla/5.0 (Linux; Android 6.0; Nexus 5 Build/MRA58N) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/136.0.0.0 Mobile Safari/537.36",
        "x-rap-xsrf-token": xsrf_token,
        "cookie": cookie_str
    }

    payload = {
        "reportId": "527b536c-e3ab-44ba-918a-9bc7fb639145",
        "pageIds": [
            "p_elblpxchtd", "p_fl86z34isd", "p_lfee22gzrd", "p_bh6r1rhhtd",
            "p_muuxu3chtd", "p_suznq25gsd", "p_zrcvvdihtd", "p_15lmn5chtd",
            "p_us2mw77gsd", "p_ykr8xr8gsd", "p_626t62ihtd", "p_nyb2wbfitd",
            "p_addj7ygitd", "p_mygp934isd", "p_regvu4gitd", "p_cf9rbohitd",
            "p_c1y6k8hitd", "p_ff7j5miitd"
        ],
        "pageNames": [
            "OPEN SLIDE", "YTD MoM Performance", "Monthly Performance Overview",
            "Monthly Performance Details", "OFFSITE SLIDE", "Facebook CPAS - Shopee",
            "Facebook CPAS - Lazada", "ONSITE SLIDE", "Lazada Sponsored Discovery",
            "Lazada Sponsored Discovery - Keyword", "Lazada Sponsored Affiliate",
            "Lazada Sponsored Store", "Lazada Sponsored Max", "Shopee Ads",
            "Shopee Ads - Keywords", "Shopee Live Ads", "Shopee Affiliate (AMS)",
            "Shopee Search Brand Ads"
        ],
        "reportName": "Marketing_performance_testing",
        "printBackground": True,
        "pdfPassword": "",
        "addReportLink": False,
        "width": 1200,
        "height": 900,
        "reportState": {
            "stateDeltas": [
                {
                    "interactions": [
                        {
                            "behaviorType": "onSelect",
                            "filterParameterValue": {
                                "filterDefinition": {
                                    "filterExpression": {
                                        "include": True,
                                        "conceptType": 0,
                                        "concept": {"name": "qt_ejrmavxtrd", "ns": "t0"},
                                        "queryTimeTransformation": {
                                            "dataTransformation": {"sourceFieldName": "_brand_name_"}
                                        },
                                        "filterConditionType": "IN",
                                        "stringValues": [brand_name]
                                    }
                                }
                            }
                        }
                    ],
                    "componentId": "cd-c7bvegnbsd"
                },
                {
                    "interactions": [
                        {
                            "behaviorType": "onSelect",
                            "filterParameterValue": {
                                "filterDefinition": {
                                    "filterExpression": {
                                        "include": True,
                                        "conceptType": 0,
                                        "concept": {"name": "qt_i7rsf6ytrd", "ns": "t0"},
                                        "queryTimeTransformation": {
                                            "dataTransformation": {"sourceFieldName": "_country_code_"}
                                        },
                                        "filterConditionType": "IN",
                                        "stringValues": [country_code]
                                    }
                                }
                            }
                        }
                    ],
                    "componentId": "cd-d2fvegnbsd"
                },
                {
                    "interactions": [
                        {
                            "behaviorType": "onSelect",
                            "dateParameterValue": {
                                "dateRange": {
                                    "startAnchor": start_date,
                                    "endAnchor": end_date,
                                    "dX": "cd-augvegnbsd"
                                }
                            }
                        }
                    ],
                    "componentId": "cd-augvegnbsd"
                }
            ],
            "pageId": "",
            "datasourceParameterOverrides": []
        },
        "pageSettings": [
            {"pageId": "p_elblpxchtd", "name": "OPEN SLIDE", "width": 2000, "height": 1125},
            {"pageId": "p_fl86z34isd", "name": "YTD MoM Performance", "width": 2000, "height": 1125},
            {"pageId": "p_lfee22gzrd", "name": "Monthly Performance Overview", "width": 2000, "height": 1125},
            {"pageId": "p_bh6r1rhhtd", "name": "Monthly Performance Details", "width": 2000, "height": 1125},
            {"pageId": "p_muuxu3chtd", "name": "OFFSITE SLIDE", "width": 2000, "height": 1125},
            {"pageId": "p_suznq25gsd", "name": "Facebook CPAS - Shopee", "width": 2000, "height": 1125},
            {"pageId": "p_zrcvvdihtd", "name": "Facebook CPAS - Lazada", "width": 2000, "height": 1125},
            {"pageId": "p_15lmn5chtd", "name": "ONSITE SLIDE", "width": 2000, "height": 1125},
            {"pageId": "p_us2mw77gsd", "name": "Lazada Sponsored Discovery", "width": 2000, "height": 1125},
            {"pageId": "p_ykr8xr8gsd", "name": "Lazada Sponsored Discovery - Keyword", "width": 2000, "height": 1125},
            {"pageId": "p_626t62ihtd", "name": "Lazada Sponsored Affiliate", "width": 2000, "height": 1125},
            {"pageId": "p_nyb2wbfitd", "name": "Lazada Sponsored Store", "width": 2000, "height": 1125},
            {"pageId": "p_addj7ygitd", "name": "Lazada Sponsored Max", "width": 2000, "height": 1125},
            {"pageId": "p_mygp934isd", "name": "Shopee Ads", "width": 2000, "height": 1125},
            {"pageId": "p_regvu4gitd", "name": "Shopee Ads - Keywords", "width": 2000, "height": 1125},
            {"pageId": "p_cf9rbohitd", "name": "Shopee Live Ads", "width": 2000, "height": 1125},
            {"pageId": "p_c1y6k8hitd", "name": "Shopee Affiliate (AMS)", "width": 2000, "height": 1125},
            {"pageId": "p_ff7j5miitd", "name": "Shopee Search Brand Ads", "width": 2000, "height": 1125}
        ]
    }

    try:
        response = requests.post(url, headers=headers, json=payload)
        if response.status_code == 200:
            try:
                pdf_data = base64.b64decode(response.content)
                print(f"Successfully downloaded PDF for {brand_name} - {country_code}")
                return pdf_data
            except Exception as decode_error:
                print(f"Error decoding PDF content: {decode_error}")
        else:
            print(f"Request failed with status code: {response.status_code}")
    except Exception as req_error:
        print(f"Error making request: {req_error}")

    return None

def validate_folder_id(folder_id: str) -> bool:
    """
    Validate if the folder ID format is correct
    """
    if not folder_id:
        return False
    # Google Drive folder IDs are typically 33 characters long
    # and contain letters, numbers, hyphens, and underscores
    return len(folder_id) >= 25 and all(c.isalnum() or c in '-_' for c in folder_id)

def get_venture_info_from_db():
    """
    Get venture information from database including parent_drive_folder_id
    Returns: List of dicts with venture, brand_name, and parent_drive_folder_id
    """
    conn = mysql.connector.connect(
        host='bi-dwh-starrocks.powersell.net',
        user='thanh.ly',
        database='media',
        port=9030,
        password='thanhALda12Dhaslznx'
    )
    cursor = conn.cursor(dictionary=True)
    
    cursor.execute("""
        SELECT DISTINCT venture, brand_name, parent_drive_folder_id 
        FROM dwh_media.dim_automation_deck_config 
        WHERE venture = 'TH' AND brand_name = 'TEFAL'
        LIMIT 2
    """)
        # WHERE parent_drive_folder_id IS NOT NULL
    rows = cursor.fetchall()
    
    # Validate folder IDs
    valid_rows = []
    for row in rows:
        folder_id = row['parent_drive_folder_id']
        if validate_folder_id(folder_id):
            valid_rows.append(row)
        else:
            print(f"Warning: Invalid folder ID for venture {row['venture']}: {folder_id}")
    
    cursor.close()
    conn.close()
    return valid_rows

def init_google_drive_service(google_auth_file):
    """
    Initialize Google Drive service using service account credentials
    """
    SCOPES = ['https://www.googleapis.com/auth/drive']
    SERVICE_ACCOUNT_FILE = google_auth_file
    
    try:
        credentials = service_account.Credentials.from_service_account_file(
            SERVICE_ACCOUNT_FILE, scopes=SCOPES)
        
        # Get and display service account email
        with open(SERVICE_ACCOUNT_FILE, 'r') as f:
            sa_info = json.load(f)
            service_account_email = sa_info.get('client_email')
            print(f"\nService Account Email: {service_account_email}")
            print("Please ensure the folder is shared with this email with Editor permission")
        
        service = build('drive', 'v3', credentials=credentials)
        return service
    except Exception as e:
        print(f"Error initializing Google Drive service: {str(e)}")
        raise

def delete_files_with_name(drive_service, folder_id: str, filename: str) -> list:
    """
    Delete all files with the same name in the specified folder
    Args:
        drive_service: Google Drive service object
        folder_id: ID of the folder to search in
        filename: Name of the file to delete
    Returns:
        list: List of deleted file IDs
    """
    deleted_files = []
    try:
        # Find all files with the same name in the folder
        query = f"'{folder_id}' in parents and name = '{filename}' and trashed = false"
        files = drive_service.files().list(
            q=query,
            fields="files(id, name)"
        ).execute().get('files', [])

        if files:
            print(f"Found {len(files)} existing files with name: {filename}")
            for file in files:
                print(f"Deleting file: {file['name']} (ID: {file['id']})")
                drive_service.files().delete(fileId=file['id']).execute()
                deleted_files.append(file['id'])
        else:
            print(f"No existing files found with name: {filename}")

    except Exception as e:
        print(f"Error while deleting files: {str(e)}")
    
    return deleted_files

def upload_file_to_drive(drive_service, venture: str, parent_folder_id: str, files_to_upload: str) -> dict:
    """
    Upload file to Google Drive with detailed error handling
    Args:
        drive_service: Google Drive service object
        venture: Venture name for logging
        parent_folder_id: ID of parent folder to upload to
        files_to_upload: Path to file to upload
    Returns:
        dict: Result of upload operation
    """
    result = {
        'uploaded_files': [],
        'failed_files': [],
        'error_message': ''
    }

    try:
        # Confirm parent folder exists
        print(f"\n=== Starting upload for venture: {venture} ===")
        print(f"Parent Folder ID: {parent_folder_id}")
        
        # Verify parent folder exists
        folder = drive_service.files().get(fileId=parent_folder_id, fields="id, name").execute()
        print(f"Confirmed parent folder: {folder.get('name')}")

        # Get file name and verify file exists
        filename = os.path.basename(files_to_upload)
        print(f"\nChecking file: {filename}")
        
        if not os.path.exists(files_to_upload):
            result['failed_files'].append({
                'path': files_to_upload,
                'error': 'File does not exist'
            })
            print(f"Error: File not found at {files_to_upload}")
            return result


        # First, delete any existing files with the same name
        deleted_files = delete_files_with_name(drive_service, parent_folder_id, filename)
        if deleted_files:
            print(f"Deleted {len(deleted_files)} existing files before upload")

        # Upload file to venture folder
        print(f"\nUploading file: {filename}")
        file_metadata = {
            'name': filename,
            'parents': [parent_folder_id]
        }
        media = MediaFileUpload(
            files_to_upload,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            resumable=True
        )
        
        file = drive_service.files().create(
            body=file_metadata,
            media_body=media,
            fields='id'
        ).execute()
        
        if file and 'id' in file:
            result['uploaded_files'].append({
                'filename': filename,
                'file_id': file.get('id'),
                'folder_id': parent_folder_id
            })
            print(f"Successfully uploaded: {filename} (ID: {file.get('id')}) to folder {venture}")
        else:
            result['failed_files'].append({
                'path': files_to_upload,
                'error': 'Upload failed - Invalid file ID'
            })
            print(f"Error: Upload failed for {filename}")

        print("\n=== Upload completed ===")
        return result

    except Exception as e:
        result['error_message'] = str(e)
        print(f"\nError during upload: {str(e)}")
        return result
    """
    Upload file to Google Drive with detailed error handling
    Args:
        drive_service: Google Drive service object
        venture: Venture name for logging
        parent_folder_id: ID of parent folder to upload to
        files_to_upload: Path to file to upload
    Returns:
        dict: Result of upload operation
    """
    result = {
        'uploaded_files': [],
        'failed_files': [],
        'error_message': ''
    }

    try:
        # Confirm parent folder exists
        print(f"\n=== Starting upload for venture: {venture} ===")
        print(f"Parent Folder ID: {parent_folder_id}")
        
        # Verify parent folder exists
        folder = drive_service.files().get(fileId=parent_folder_id, fields="id, name").execute()
        print(f"Confirmed parent folder: {folder.get('name')}")

        # Get file name and verify file exists
        filename = os.path.basename(files_to_upload)
        print(f"\nChecking file: {filename}")
        
        if not os.path.exists(files_to_upload):
            result['failed_files'].append({
                'path': files_to_upload,
                'error': 'File does not exist'
            })
            print(f"Error: File not found at {files_to_upload}")
            return result

        # Upload file directly without checking for duplicates
        print(f"\nUploading file: {filename}")
        file_id = upload_file(
            file_obj=files_to_upload,
            filename=filename,
            parent_folder_id=parent_folder_id,
            file_type='pptx'
        )
        
        if file_id:
            result['uploaded_files'].append({
                'filename': filename,
                'file_id': file_id
            })
            print(f"Successfully uploaded: {filename} (ID: {file_id})")
        else:
            result['failed_files'].append({
                'path': files_to_upload,
                'error': 'Upload failed - Invalid file ID'
            })
            print(f"Error: Upload failed for {filename}")

        print("\n=== Upload completed ===")
        return result

    except Exception as e:
        result['error_message'] = str(e)
        print(f"\nError during upload: {str(e)}")
        return result

def process_batch_reports(reports_config: List[Dict], base_output_dir: str, cookies_file: str, start_date:str, end_date:str, scale: float = 1.0) -> List[Dict]:
    """
    Process multiple reports based on the provided configuration.
    Args:
        reports_config: List of dictionaries containing brand_name and venture
        base_output_dir: Base directory for all outputs
        cookies_file: Path to cookies.json
        scale: Scale factor for image size in PPTX
    Returns:
        List of dictionaries containing results for each report
    """
    results = []
    
    for report in reports_config:
        brand_name = report['brand_name']
        venture = report['venture']
        print(f"\nProcessing report for {brand_name} - {venture}")
        # Create venture-specific directory
        venture_dir = os.path.join(base_output_dir, venture)
        os.makedirs(venture_dir, exist_ok=True)
        # Format dates as YYYYMMDD
        start_date_str = start_date.replace('-', '')
        end_date_str = end_date.replace('-', '')
        pptx_path = os.path.join(venture_dir, f"{brand_name}_{start_date_str}_{end_date_str}.pptx")
        # Download PDF as bytes
        pdf_bytes = download_looker_report_bytes(
            brand_name=brand_name,
            country_code=venture,
            start_date=start_date,
            end_date=end_date,  # Using venture as country_code for Looker filter
            cookies_file_path=cookies_file
        )
        result = {
            'brand_name': brand_name,
            'venture': venture,
            'pptx_path': '',
            'status': 'failed'
        }
        if pdf_bytes:
            pptx_result = pdf_to_pptx_from_bytes(pdf_bytes, pptx_path, scale=scale)
            if pptx_result:
                result['pptx_path'] = pptx_result
                result['status'] = 'success'
                print(f"Successfully processed report for {brand_name} - {venture}")
            else:
                print(f"Failed to convert PDF to PPTX for {brand_name} - {venture}")
        else:
            print(f"Failed to download report for {brand_name} - {venture}")
        results.append(result)
        time.sleep(2)
    return results