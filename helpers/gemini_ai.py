import os
import io
import sys
from PIL import Image
from pptx.dml.color import RGBColor
from pptx import Presentation
from PIL import Image
from  google import generativeai as genai
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.enum.text import MSO_AUTO_SIZE


def summarize_pptx_with_gemini(pptx_path, skip_slides=[]):
    api_key = 'key'

    # Configure Gemini API
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel(model_name="gemini-2.5-flash-preview-05-20")


    # Load PowerPoint
    prs = Presentation(pptx_path)
    story_so_far = ""  # To maintain continuity
    summaries = []

    # Process each slide
    for i, slide in enumerate(prs.slides):
        
        # Skip Open slides
        if i+1 in skip_slides:
            continue
               
        # Extract first image (if any)
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                break  # Use only one image per slide for simplicity

        image_new = Image.open(io.BytesIO(image_bytes))
        
        # Build prompt
        prompt = f"""
        This is slide {i+1} of a presentation. 
        The slide contains an image that may provide context for the summary.
        The general context of the presentation is about monthly marketing performance analysis for a brand.
        The story so far: "{story_so_far}".
        Please summarize this slide's content in a way that continues the narrative logically and smoothly based on the slide with no more than 50 words.
        Try not to repeat information already provided in the story so far as well as the way of summarizing the previous slides.
        The summary should be concise and informative, focusing on the key points of the slide.
        Remember to focus on the current month's marketing performance and how it relates to the overall brand strategy.
        The tone should be professional and suitable for a business presentation.
"""

        # Generate response
        response = model.generate_content(
            contents=[
                    prompt,
                    image_new
            ]
        )

        summary = response.text.lstrip('\n\r')

        left = Inches(0.3)
        top = Inches(0.5)
        width = Inches(9.37)
        height = Inches(0.7)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        
        fill = txBox.fill
        fill.solid()  # Set the fill to a solid color
        fill.fore_color.rgb = RGBColor(236, 240, 254)  # Set to white
        
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        tf.text = summary  # Add directly to remove the empty first line
        p = tf.paragraphs[0]
        p.font.size = Pt(12)  # Optional: Skip this if auto_size should control it

        # Optional: Set font size, color, etc. 

        # Store result and update story
        # summaries.append(summary)
        story_so_far += f" Slide {i+1}: {summary}"
    
    # output_pptx_path = pptx_path.replace(".pptx", "_with_summaries.pptx")
    prs.save(pptx_path)
    print(f"Modified presentation saved to: {pptx_path}")

    return summaries

